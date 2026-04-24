"""
PDF and document text extraction utilities.
"""

import os
import re
from collections import Counter

import PyPDF2
import pdfplumber
from docx import Document


_URL_RE = re.compile(r"https?://\S+|www\.\S+", re.I)
_BULLET_RE = re.compile(r"^[\s\-\*>]+")
_META_RE = re.compile(
    r"^(who has|raise your hand|poll|course outline|agenda|prepared by|adapted from|"
    r"modified from|based on material|copyright|source:|credit:|image:|figure \d|"
    r"table \d|slide \d|click here|sign up|subscribe)",
    re.I,
)
_BAD_CONTAINS = [
    "with material from",
    "your institute",
    "your name",
    "who has programming experience",
]


def _normalize_line(line: str) -> str:
    line = line.strip()
    line = _URL_RE.sub("", line)
    line = _BULLET_RE.sub("", line)
    line = re.sub(r"\s+", " ", line).strip()
    return line


def _is_junk_line(line: str) -> bool:
    if not line:
        return True
    low = line.lower()
    if len(low) < 3:
        return True
    if _META_RE.search(low):
        return True
    if any(bad in low for bad in _BAD_CONTAINS):
        return True
    if re.fullmatch(r"[\W_]+", line):
        return True
    return False


def clean_text(text: str) -> str:
    """Conservative cleaning that preserves useful academic content."""
    if not text:
        return ""

    clean_lines = []
    seen = set()
    for raw in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        line = _normalize_line(raw)
        if not line or _is_junk_line(line):
            continue

        # Keep lines that likely carry useful content.
        if len(line) < 5:
            continue
        if len(line.split()) <= 2 and line.isupper():
            continue

        norm = line.lower()
        if norm in seen:
            continue
        seen.add(norm)
        clean_lines.append(line)

    return "\n".join(clean_lines).strip()


def extract_text_from_pdf(file_path: str) -> str | None:
    """Extract text from PDF and remove repeated boilerplate lines."""
    try:
        with open(file_path, "rb") as fh:
            reader = PyPDF2.PdfReader(fh)
            pages_lines = []
            for page in reader.pages:
                page_text = page.extract_text() or ""
                lines = [ln.strip() for ln in page_text.splitlines() if ln.strip()]
                pages_lines.append(lines)

        page_count = max(len(pages_lines), 1)
        freq = Counter(ln for pg in pages_lines for ln in pg)
        boiler = set()
        if page_count >= 2:
            boiler = {
                ln for ln, c in freq.items()
                if len(ln) <= 90 and (c / page_count) >= 0.35
            }

        kept = [ln for pg in pages_lines for ln in pg if ln not in boiler]
        cleaned = clean_text("\n".join(kept))
        print(
            f"[PDF] {os.path.basename(file_path)} -> pages={page_count} "
            f"boiler_removed={len(boiler)} chars={len(cleaned)}"
        )
        return cleaned or None
    except Exception as e:
        print(f"Error extracting PDF: {e}")
        return None


def extract_text_from_docx(file_path: str) -> str | None:
    try:
        doc = Document(file_path)
        text = "\n".join(p.text for p in doc.paragraphs)
        cleaned = clean_text(text)
        print(f"[DOCX] {os.path.basename(file_path)} -> chars={len(cleaned)}")
        return cleaned or None
    except Exception as e:
        print(f"Error extracting DOCX: {e}")
        return None


def extract_text_from_pptx(file_path: str) -> str | None:
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        blocks = []
        for slide in prs.slides:
            raw_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    raw_lines.extend(shape.text.split("\n"))
            slide_text = clean_text("\n".join(raw_lines))
            if slide_text:
                blocks.append(slide_text)

        cleaned = "\n\n".join(blocks).strip()
        print(
            f"[PPTX] {os.path.basename(file_path)} -> "
            f"slides={len(prs.slides)} kept={len(blocks)} chars={len(cleaned)}"
        )
        return cleaned or None
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
        return None


def _extract_pdf_with_pdfplumber(file_path: str) -> str:
    """Fallback PDF extraction path when PyPDF2 output is empty."""
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n\n"
    except Exception as e:
        print(f"pdfplumber fallback failed: {e}")
    return clean_text(text)


def extract_text_from_file(file_path: str) -> str:
    """Auto-detect file type and extract text."""
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    if ext == ".pdf":
        text = extract_text_from_pdf(file_path) or ""
        if not text:
            text = _extract_pdf_with_pdfplumber(file_path) or ""
    elif ext in [".docx", ".doc"]:
        text = extract_text_from_docx(file_path) or ""
    elif ext in [".pptx", ".ppt"]:
        text = extract_text_from_pptx(file_path) or ""
    else:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
                text = clean_text(fh.read())
        except Exception:
            text = ""

    print(f"[DEBUG] extract_text_from_file -> {os.path.basename(file_path)} chars={len(text)}")
    return text
